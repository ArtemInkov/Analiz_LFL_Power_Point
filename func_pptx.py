from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.chart.data import CategoryChartData
from pptx.util import Inches
from pptx.dml.color import RGBColor
import pptx.enum.chart
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt
from openpyxl import load_workbook
import aspose.slides as slides
import aspose.pydrawing as drawing
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import  MSO_AUTO_SIZE
import xlrd
a = 'анализ продаж.xlsx'
b = 'Анализ продаж по кварталам.pptx'
def presentation_pptx_ferst(file_name = ''):
    """Функция создает файл продажи и первый слайд берет года из файла exel"""
    file_xls = xlrd.open_workbook(file_name)
    sheets_xls = file_xls.sheet_by_name("Анализ год к году")
    last_year = sheets_xls.cell(1, 1).value.split()
    last_year = last_year[1]
    year = sheets_xls.cell(1, 2).value.split()
    year = year[1]
    presentation = Presentation()
    print(presentation.slides)
    ferst_slide = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(ferst_slide)
    slide.shapes.title.text = f'Продажи к сравнению\n' \
                              f'{last_year} - {year}г.'
    slide.shapes.title.width = Inches(14)
    presentation.save('Анализ продаж по кварталам.pptx')



def presentation_pptx_shops_graphic(file_name_xls = '', file_name_pptx = ''):
    file_xls = xlrd.open_workbook(file_name_xls)
    sheets_xls = file_xls.sheet_by_name("Анализ год к году")
    colum_sh = sheets_xls.ncols
    row_sh = sheets_xls.nrows
    last_year = sheets_xls.cell(1, 1).value.split()
    last_year = last_year[1]
    year = sheets_xls.cell(1, 2).value.split()
    year = year[1]
    p = Presentation(file_name_pptx)

    print(range(int(colum_sh - 1)//4))
    for i in range(1, int((colum_sh - 1)//4)+ 1):
        print(int((colum_sh - 1)//4))
        print(i)
        if i == 1:
            mount_title = sheets_xls.cell(1, 1).value.split()
        else:
            mount_title = sheets_xls.cell(1, (i * 4)-3).value.split()
        print(mount_title)
        p.slide_width = Inches(15)
        ferst_slide = p.slide_layouts[5]

        slide = p.slides.add_slide(ferst_slide)
        s = p.slide_width
        print(s)
        slide.shapes.title.text = f'\nПродажи к сравнению в рублях\n{mount_title[0]}'
        slide.shapes.title.width = Inches(14)
        chart_data = ChartData()
        shop_list = []
        last_year_value = []
        current_year_value = []
        procent_year = []
        summ_last_year_rub = []
        summ_current_year_rub = []
        summ_last_year_weight = []
        summ_current_year_weight = []

        for i_shop in range(1, row_sh):
            value_cell = sheets_xls.cell(i_shop, 0).value
            if value_cell != '' and value_cell != 'Магазин':
                shop_list.append(sheets_xls.cell(i_shop, 0).value)
                print(sheets_xls.cell(i_shop, 0).value)
                if i == 1:
                    #print(sheets_xls.cell(i_shop, 1).value)
                    #print(sheets_xls.cell(i_shop, 2).value)
                    if isinstance(sheets_xls.cell(i_shop, 1).value, float):
                        last_year_value.append(sheets_xls.cell(i_shop, 1).value)
                    else:
                        last_year_value.append(0)
                    if isinstance(sheets_xls.cell(i_shop, 2).value, float):
                        current_year_value.append(sheets_xls.cell(i_shop, 2).value)
                    else:
                        current_year_value.append(0)
                    if isinstance(sheets_xls.cell(i_shop, 1).value, float) and isinstance(sheets_xls.cell(i_shop, 2).value, float):
                        procent_year.append(round((100 - (sheets_xls.cell(i_shop, 1).value / sheets_xls.cell(i_shop, 2).value * 100)), 2))
                    else:
                        procent_year.append(0)
                else:
                    a1 = sheets_xls.cell(i_shop, (i * 4)-3).value
                    a = sheets_xls.cell(i_shop, (i * 4)-2).value
                    if isinstance(a1, float):
                        last_year_value.append(sheets_xls.cell(i_shop, (i * 4)-3).value)
                    else:
                        last_year_value.append(0)
                    if isinstance(a, float):
                        current_year_value.append(sheets_xls.cell(i_shop, (i * 4)-2).value)
                    else:
                        current_year_value.append(0)
                    if isinstance(a1, float) and isinstance(a, float):
                        procent_year.append(round(100 - (sheets_xls.cell(i_shop, (i * 4)-3).value / sheets_xls.cell(i_shop, (i * 4)-2).value * 100), 2))
                    else:
                        procent_year.append(0)
            print(shop_list)
            print(last_year_value)
            print(current_year_value)
            print(procent_year)

        chart_data.categories = shop_list
        chart_data.add_series(last_year, (tuple(last_year_value)))
        chart_data.add_series(year, (tuple(current_year_value)))
        chart_data.add_series('%', (tuple(procent_year)))

        summ_last_year_rub.append(sum(last_year_value))
        summ_current_year_rub.append(sum(current_year_value))



        x, y, cx, cy = Inches(0.01), Inches(2), Inches(15), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
        chart = graphic_frame.chart
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels

        data_labels.font.size = Pt(13)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

        """
        Создание слайда с продажами вес
        
        """

        p.slide_width = Inches(15)
        ferst_slide = p.slide_layouts[5]

        # help(ferst_slide)
        slide = p.slides.add_slide(ferst_slide)
        s = p.slide_width
        print(s)
        slide.shapes.title.text = f'\nПродажи к сравнению вес\n{mount_title[0]}'
        slide.shapes.title.width = Inches(14)
        chart_data = ChartData()
        shop_list = []
        last_year_value = []
        current_year_value = []
        procent_year = []
        for i_shop in range(3, row_sh):
            value_cell = sheets_xls.cell(i_shop, 0).value
            if value_cell != '' and value_cell != 'Магазин':
                shop_list.append(sheets_xls.cell(i_shop, 0).value)
                print(sheets_xls.cell(i_shop, 0).value)
                if i == 3:
                    print(sheets_xls.cell(i_shop, 3).value)
                    print(sheets_xls.cell(i_shop, 4).value)

                    if isinstance(sheets_xls.cell(i_shop, 3).value, float):
                        last_year_value.append(sheets_xls.cell(i_shop, 3).value)
                    else:
                        last_year_value.append(0)
                    if isinstance (sheets_xls.cell(i_shop, 4).value, float):
                        current_year_value.append(sheets_xls.cell(i_shop, 4).value)
                    else:
                        current_year_value.append(0)
                    if isinstance(sheets_xls.cell(i_shop, 3).value, float) and isinstance(
                            sheets_xls.cell(i_shop, 4).value, float):
                        procent_year.append(
                            round((100 - (sheets_xls.cell(i_shop, 3).value / sheets_xls.cell(i_shop, 4).value * 100)),
                                  2))
                    else:
                        procent_year.append(0)
                else:
                    a1 = sheets_xls.cell(i_shop, (i * 4) - 1).value
                    a = sheets_xls.cell(i_shop, (i * 4)).value
                    if isinstance(a1, float):
                        last_year_value.append(sheets_xls.cell(i_shop, (i * 4) - 1).value)
                    else:
                        last_year_value.append(0)
                    if isinstance(a, float):
                        current_year_value.append(sheets_xls.cell(i_shop, (i * 4)).value)
                    else:
                        current_year_value.append(0)
                    if isinstance(a1, float) and isinstance(a, float):
                        procent_year.append(round(100 - (
                                    sheets_xls.cell(i_shop, (i * 4) - 1).value / sheets_xls.cell(i_shop, (
                                        i * 4)).value * 100), 2))
                    else:
                        procent_year.append(0)
            print(shop_list)
            print(last_year_value)
            print(current_year_value)
            print(procent_year)

        chart_data.categories = shop_list
        chart_data.add_series(last_year, (tuple(last_year_value)))
        chart_data.add_series(year, (tuple(current_year_value)))
        chart_data.add_series('%', (tuple(procent_year)))

        summ_last_year_weight.append(sum(last_year_value))
        summ_current_year_weight.append(sum(current_year_value))

        x, y, cx, cy = Inches(0.01), Inches(2), Inches(15), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
        chart = graphic_frame.chart
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels

        data_labels.font.size = Pt(13)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        """
        Создание слайда с общими показателями за месяц
        """
        p.slide_width = Inches(15)
        ferst_slide = p.slide_layouts[5]

        slide = p.slides.add_slide(ferst_slide)
        s = p.slide_width
        print(s)
        slide.shapes.title.text = f'\n                Продажи к сравнению                  \nв киллограмах                               в рублях'

        slide.shapes.title.width = Inches(14)
        chart_data = ChartData()

        summ_procent_year_rub = []
        procent = round(100 - summ_last_year_rub[0] / summ_current_year_rub[0] * 100, 1)
        summ_procent_year_rub.append(procent)

        chart_data.categories = ['в рублях']
        chart_data.add_series(last_year, tuple(summ_last_year_rub))
        chart_data.add_series(year, tuple(summ_current_year_rub))
        chart_data.add_series('%', tuple(summ_procent_year_rub))



        x, y, cx, cy = Inches(1), Inches(2), Inches(6), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
        chart = graphic_frame.chart
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels

        data_labels.font.size = Pt(13)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

        chart_data_weight = ChartData()

        summ_procent_year_weight = []
        procent = round(100 - summ_last_year_weight[0] / summ_current_year_weight[0] * 100, 1)
        summ_procent_year_weight.append(procent)

        chart_data_weight.categories = ['В киллограмах']
        chart_data_weight.add_series(last_year, tuple(summ_last_year_weight))
        chart_data_weight.add_series(year, tuple(summ_current_year_weight))
        chart_data_weight.add_series('%', tuple(summ_procent_year_weight))

        x, y, cx, cy = Inches(9), Inches(2), Inches(6), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data_weight)
        chart = graphic_frame.chart
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels

        data_labels.font.size = Pt(13)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

        p.save('chart-01.pptx')

def presentation_pptx_shops_graphic_quarter(file_name_xls = '', file_name_pptx = ''):
    file_xls = xlrd.open_workbook(file_name_xls)
    sheets_xls = file_xls.sheet_by_name("Анализ год к году")
    colum_sh = sheets_xls.ncols
    row_sh = sheets_xls.nrows
    last_year = sheets_xls.cell(1, 1).value.split()
    last_year = last_year[1]
    year = sheets_xls.cell(1, 2).value.split()
    year = year[1]
    p = Presentation(file_name_pptx)
    quarter_mounth = []
    count_mounth = 0
    for i in range(1, int((colum_sh - 1)//4) // 3 + 1):
        print(int((colum_sh - 1)//4) / 3)
        print(i)
        shop_list = []
        last_year_value_rub = []
        current_year_value_rub = []
        last_year_value_weight = []
        current_year_value_weight = []
        procent_year_rub = []
        procent_year_weight = []
        summ_last_year_rub = []
        summ_current_year_rub = []
        summ_last_year_weight = []
        summ_current_year_weight = []
        summ_last_year_quarter_rub = 1
        summ_current_year_quarter_rub = 1
        summ_last_year_quarter_weight = 1
        summ_current_year_quarter_weight = 1

        for i_shop in range(1, row_sh):
            value_cell = sheets_xls.cell(i_shop, 0).value
            if value_cell != '' and value_cell != 'Магазин':
                shop_list.append(sheets_xls.cell(i_shop, 0).value)
                print(sheets_xls.cell(i_shop, 0).value)
                if i == 1:
                    #print(sheets_xls.cell(i_shop, 1).value)
                    #print(sheets_xls.cell(i_shop, 2).value)
                    for i_value in range(1, 4):
                        if isinstance(sheets_xls.cell(i_shop, i_value * 4).value, float):
                            summ_last_year_quarter_rub += (sheets_xls.cell(i_shop, i_value * 4).value)
                        if isinstance(sheets_xls.cell(i_shop, (i_value + 1) * 4).value, float):
                            print((sheets_xls.cell(i_shop,(i_value + 1) * 4).value))
                            summ_current_year_quarter_rub += (sheets_xls.cell(i_shop,(i_value + 1) * 4).value)
                        if isinstance(sheets_xls.cell(i_shop, (i_value + 2) * 4).value, float):
                            summ_last_year_quarter_weight += (sheets_xls.cell(i_shop,(i_value + 2) * 4).value)
                        if isinstance(sheets_xls.cell(i_shop, (i_value + 3) * 4).value, float):
                            summ_current_year_quarter_weight += (sheets_xls.cell(i_shop,(i_value + 3) * 4).value)
                # else:
                #     a1 = sheets_xls.cell(i_shop, (i * 4)-3).value
                #     a = sheets_xls.cell(i_shop, (i * 4)-2).value
                #     if isinstance(a1, float):
                #         last_year_value.append(sheets_xls.cell(i_shop, (i * 4)-3).value)
                #     else:
                #         last_year_value.append(0)
                #     if isinstance(a, float):
                #         current_year_value.append(sheets_xls.cell(i_shop, (i * 4)-2).value)
                #     else:
                #         current_year_value.append(0)
                #     if isinstance(a1, float) and isinstance(a, float):
                #         procent_year.append(round(100 - (sheets_xls.cell(i_shop, (i * 4)-3).value / sheets_xls.cell(i_shop, (i * 4)-2).value * 100), 2))
                #     else:
                #         procent_year.append(0)
                last_year_value_rub.append(summ_last_year_quarter_rub)
                current_year_value_rub.append(summ_current_year_quarter_rub)
                last_year_value_weight.append(summ_last_year_quarter_weight)
                current_year_value_weight.append(summ_current_year_quarter_weight)
                procent_year_rub.append(round((1 - summ_last_year_quarter_rub / summ_current_year_quarter_rub) * 100, 2))
                procent_year_weight.append(round((1 - summ_last_year_quarter_weight / summ_current_year_quarter_weight ) * 100 , 2))
                print(last_year_value_weight)
                summ_last_year_rub.append(sum(last_year_value_rub))
                summ_current_year_rub.append(sum(current_year_value_rub))
                summ_last_year_weight.append(sum(last_year_value_weight))
                summ_current_year_weight.append(sum(current_year_value_weight))
                summ_last_year_quarter_rub = 1
                summ_current_year_quarter_rub = 1
                summ_last_year_quarter_weight = 1
                summ_current_year_quarter_weight = 1

        """Создал первый слайд, разобраться с процентами"""
        p.slide_width = Inches(15)
        ferst_slide = p.slide_layouts[5]
        slide = p.slides.add_slide(ferst_slide)
        s = p.slide_width
        slide.shapes.title.text = f'\nПродажи к сравнению в рублях 1 квартал \n()'
        slide.shapes.title.width = Inches(14)
        chart_data = ChartData()
        chart_data.categories = shop_list
        chart_data.add_series(last_year, (tuple(last_year_value_rub)))
        chart_data.add_series(year, (tuple(current_year_value_rub)))
        chart_data.add_series('%', (tuple(procent_year_rub)))

        x, y, cx, cy = Inches(0.01), Inches(2), Inches(15), Inches(4.5)
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
        chart = graphic_frame.chart
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels

        data_labels.font.size = Pt(13)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    """Скопировать создание первого слайда, передалть под вес
    потом слайд общий """
    p.save('chart-01.pptx')


presentation_pptx_ferst(a)
presentation_pptx_shops_graphic(a, b)
presentation_pptx_shops_graphic_quarter(a, b)

